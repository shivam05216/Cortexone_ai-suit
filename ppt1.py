from pptx import Presentation

# Create a new presentation for Problem Statement ID 1624
new_electricity_ppt = Presentation()

# Slide 1: Title Page
slide_1_layout = new_electricity_ppt.slide_layouts[0]  # Title slide layout
slide_1 = new_electricity_ppt.slides.add_slide(slide_1_layout)

title_placeholder = slide_1.shapes.title
title_placeholder.text = "SMART INDIA HACKATHON 2024\nAI-based Model for Electricity Demand Projection in Delhi"

subtitle_placeholder = slide_1.placeholders[1]
subtitle_placeholder.text = "Problem Statement ID: 1624\nTheme: Smart Automation\nCategory: Software\nTeam Name: (Your Team Name)\nTeam ID: (Your Team ID)"

# Slide 2: Proposed Solution
slide_2_layout = new_electricity_ppt.slide_layouts[1]  # Title and content layout
slide_2 = new_electricity_ppt.slides.add_slide(slide_2_layout)

title_2 = slide_2.shapes.title
title_2.text = "Proposed Solution"

content_2 = slide_2.shapes[1]
content_2.text = """The AI-based model will predict electricity demand and peak load by considering key factors like weather, real estate development, and public holidays.
Key Features:
1. AI-based forecasting using historical data and real-time factors (weather, holidays).
2. Weather compensation: Factoring in temperature, humidity, wind speed, and rainfall.
3. Public holiday and weekday analysis for accurate peak predictions.
4. Dual-peak load prediction during summer months.
5. Dynamic adjustment based on real estate growth and uneven load growth in Delhi.
6. Compensation for solar generation and Duck Curve effects."""

# Slide 3: Technical Approach
slide_3 = new_electricity_ppt.slides.add_slide(slide_2_layout)

title_3 = slide_3.shapes.title
title_3.text = "Technical Approach"

content_3 = slide_3.shapes[1]
content_3.text = """Technologies:
- Machine learning models for demand forecasting.
- Integration with real-time weather APIs.
- AI-driven adjustments for peak loads and solar compensation.
- Cloud-based infrastructure for real-time data processing.

Methodology:
1. Data Collection: Historical consumption data, weather data, real estate data.
2. AI-Driven Forecasting: Machine learning to predict demand trends.
3. Peak Load Identification: Analyze patterns to predict summer/winter peak loads.
4. Solar Compensation: Account for Duck Curve and CERC variations in solar energy."""

# Slide 4: Feasibility & Viability
slide_4 = new_electricity_ppt.slides.add_slide(slide_2_layout)

title_4 = slide_4.shapes.title
title_4.text = "Feasibility & Viability"

content_4 = slide_4.shapes[1]
content_4.text = """Challenges:
1. Large data volume for processing.
2. Real-time weather adjustments.
3. Handling high variation in demand across areas.

Solutions:
1. Scalable cloud infrastructure for data storage and real-time processing.
2. Efficient AI algorithms for weather and peak load adjustments.
3. Regular model training using dynamic, real-time data."""

# Slide 5: Impact & Benefits
slide_5 = new_electricity_ppt.slides.add_slide(slide_2_layout)

title_5 = slide_5.shapes.title
title_5.text = "Impact & Benefits"

content_5 = slide_5.shapes[1]
content_5.text = """Impact:
1. Accurate forecasting of electricity demand and peak loads.
2. Better planning and distribution of electricity.
3. Reduced risk of over-purchase or under-supply of electricity.

Benefits:
1. Economic: Optimized power purchase and reduced operational costs.
2. Environmental: Improved integration of renewable energy sources.
3. Social: Reliable electricity supply, reducing blackouts and fluctuations."""

# Slide 6: Research & References
slide_6 = new_electricity_ppt.slides.add_slide(slide_2_layout)

title_6 = slide_6.shapes.title
title_6.text = "Research & References"

content_6 = slide_6.shapes[1]
content_6.text = """References:
1. Delhi Power Grid Reports.
2. AI-based demand forecasting: Current research papers.
3. Real-time data integration techniques."""

# Save the new AI Electricity presentation
new_electricity_ppt_path = "/mnt/data/AI_Electricity_Demand_SIH_2024.pptx"
new_electricity_ppt.save(new_electricity_ppt_path)
